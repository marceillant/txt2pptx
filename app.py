from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

#9/14/2023

def create_pptx_from_txt(filename):
    # Create a blank presentation
    presentation = Presentation()

    # Open the text file
    with open(filename, 'r') as file:
        lines = file.readlines()

    # Split the lines into groups of 4
    groups = [lines[n:n+4] for n in range(0, len(lines), 4)]

    # For each group of 4 lines, create a new slide with a title and content
    for group in groups:
        slide_layout = presentation.slide_layouts[6]  # Use the 'title slide' layout
        slide = presentation.slides.add_slide(slide_layout)

        # Set the slide background color to black
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black color

        # Calculate the position and size of the text box
        left = Inches(1)
        top = Inches(1)
        width = presentation.slide_width - 2 * Inches(1)
        height = presentation.slide_height - 2 * Inches(1)

        # Add a text box to the slide and center the text
        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox.fill.background()
        tf = txBox.text_frame

        # Clear all paragraphs in the text frame
        tf.clear()  # This is the new line of code

        # Center text vertically in the text box
        tf.text_anchor = MSO_ANCHOR.MIDDLE

        p = tf.add_paragraph()
        p.text = '\n'.join(group)
        p.alignment = PP_ALIGN.CENTER

        # Set the font size and color
        for run in p.runs:
            run.font.size = Pt(48)
            run.font.color.rgb = RGBColor(255, 255, 255)  # White color

    # Save the presentation
    presentation.save('Lyrics.pptx')
    
    

# Call the function with the text file
create_pptx_from_txt('Lyrics.txt')